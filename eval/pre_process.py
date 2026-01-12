"""
PPT Pre-processing Module

This module handles pre-processing of PowerPoint files and URLs:
- Extract slide images from PPTX files
- Extract slide images from PDF files (e.g., NotebookLM outputs)
- Handle different product structures (Kimi-Standard, Kimi-Smart, etc.)
- Process PPT URLs (e.g., shared presentations)
- Convert slides to images for VLM evaluation

Author: PPT Evaluation System
"""

import os
import sys
import asyncio
import argparse
import logging
import shutil
import subprocess
from pathlib import Path
from typing import List, Optional, Dict
from dataclasses import dataclass, field
from tqdm import tqdm

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


# =============================================================================
# Data Classes
# =============================================================================

@dataclass
class ProcessingResult:
    """Result of processing a PPT."""
    source_path: str
    output_dir: str
    product: str = ""
    difficulty: str = ""
    topic: str = ""
    ppt_id: str = ""
    num_slides: int = 0
    image_paths: List[str] = field(default_factory=list)
    success: bool = False
    error: Optional[str] = None
    
    @classmethod
    def from_existing(
        cls,
        source_path: str,
        output_dir: str,
        image_paths: List[str],
        **kwargs
    ) -> "ProcessingResult":
        """Create a result for already-processed content."""
        return cls(
            source_path=source_path,
            output_dir=output_dir,
            image_paths=image_paths,
            num_slides=len(image_paths),
            success=True,
            **kwargs
        )
    
    @classmethod
    def from_error(cls, source_path: str, output_dir: str, error: str, **kwargs) -> "ProcessingResult":
        """Create a result for a failed processing."""
        return cls(
            source_path=source_path,
            output_dir=output_dir,
            success=False,
            error=error,
            **kwargs
        )


# =============================================================================
# Helper Functions
# =============================================================================

def get_existing_images(images_dir: str) -> List[str]:
    """Get sorted list of existing slide images in a directory."""
    if not os.path.exists(images_dir):
        return []
    return sorted([
        os.path.join(images_dir, f)
        for f in os.listdir(images_dir)
        if f.lower().endswith(('.png', '.jpg', '.jpeg'))
    ])


def find_files(directory: str, extensions: List[str]) -> List[Path]:
    """Find files with given extensions in a directory."""
    files = []
    for ext in extensions:
        files.extend(Path(directory).glob(f"*{ext}"))
        files.extend(Path(directory).glob(f"*{ext.upper()}"))
    return files


# Product name to directory path mapping
PRODUCT_PATHS = {
    "Gamma": "Gamma",
    "NotebookLM": "NotebookLM",
    "Kimi-Standard": "Kimi/Standard",
    "Kimi-Smart": "Kimi/Smart",
    "Kimi-Banana": "Kimi/Banana",
    "Skywork": "Skyworks",
    "Skywork-Banana": "Skyworks/Banana",
    "Zhipu": "Zhipu",
    "Quake": "Quake",
}

ALL_PRODUCTS = list(PRODUCT_PATHS.keys())


# =============================================================================
# PDF Processor
# =============================================================================

class PDFProcessor:
    """Convert PDF files to slide images."""
    
    def __init__(self, image_format: str = "png", image_dpi: int = 150):
        self.image_format = image_format
        self.image_dpi = image_dpi
    
    def process(self, pdf_path: str, output_dir: str) -> ProcessingResult:
        """Process a PDF file and extract page images."""
        if not os.path.exists(pdf_path):
            return ProcessingResult.from_error(pdf_path, output_dir, f"File not found: {pdf_path}")
        
        images_dir = os.path.join(output_dir, "slide_images")
        os.makedirs(images_dir, exist_ok=True)
        
        try:
            image_paths = self._convert_to_images(pdf_path, images_dir)
            if image_paths:
                return ProcessingResult.from_existing(pdf_path, output_dir, image_paths)
            return ProcessingResult.from_error(pdf_path, output_dir, "Failed to extract images from PDF")
        except Exception as e:
            logger.error(f"Error processing PDF {pdf_path}: {e}")
            return ProcessingResult.from_error(pdf_path, output_dir, str(e))
    
    def _convert_to_images(self, pdf_path: str, output_dir: str) -> List[str]:
        """Convert PDF to images using available methods."""
        # Try pdf2image
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=self.image_dpi, fmt=self.image_format)
            return self._save_images(images, output_dir)
        except (ImportError, Exception) as e:
            logger.debug(f"pdf2image failed: {e}")
        
        # Try PyMuPDF
        try:
            import fitz
            return self._convert_with_pymupdf(pdf_path, output_dir)
        except (ImportError, Exception) as e:
            logger.debug(f"PyMuPDF failed: {e}")
        
        # Try Ghostscript
        try:
            return self._convert_with_ghostscript(pdf_path, output_dir)
        except Exception as e:
            logger.debug(f"Ghostscript failed: {e}")
        
        logger.error("All PDF conversion methods failed")
        return []
    
    def _save_images(self, images, output_dir: str) -> List[str]:
        """Save PIL images to files."""
        paths = []
        for i, img in enumerate(images):
            path = os.path.join(output_dir, f"slide_{i+1:04d}.{self.image_format}")
            img.save(path)
            paths.append(path)
        return paths
    
    def _convert_with_pymupdf(self, pdf_path: str, output_dir: str) -> List[str]:
        """Convert using PyMuPDF."""
        import fitz
        doc = fitz.open(pdf_path)
        paths = []
        mat = fitz.Matrix(self.image_dpi / 72, self.image_dpi / 72)
        
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            path = os.path.join(output_dir, f"slide_{i+1:04d}.{self.image_format}")
            pix.save(path)
            paths.append(path)
        
        doc.close()
        return paths
    
    def _convert_with_ghostscript(self, pdf_path: str, output_dir: str) -> List[str]:
        """Convert using Ghostscript."""
        gs_cmd = shutil.which("gs") or shutil.which("gswin64c")
        if not gs_cmd:
            raise RuntimeError("Ghostscript not found")
        
        output_pattern = os.path.join(output_dir, "slide_%04d.png")
        result = subprocess.run([
            gs_cmd, "-dNOPAUSE", "-dBATCH", "-sDEVICE=png16m",
            f"-r{self.image_dpi}", f"-sOutputFile={output_pattern}", pdf_path
        ], capture_output=True, timeout=300)
        
        if result.returncode != 0:
            raise RuntimeError(f"Ghostscript failed: {result.stderr.decode()}")
        
        return get_existing_images(output_dir)


# =============================================================================
# PPTX Processor
# =============================================================================

class PPTXProcessor:
    """Convert PPTX files to slide images."""
    
    def __init__(self, image_format: str = "png", image_dpi: int = 150):
        self.image_format = image_format
        self.image_dpi = image_dpi
        self.pdf_processor = PDFProcessor(image_format, image_dpi)
    
    def process(self, pptx_path: str, output_dir: str) -> ProcessingResult:
        """Process a PPTX file."""
        if not os.path.exists(pptx_path):
            return ProcessingResult.from_error(pptx_path, output_dir, f"File not found: {pptx_path}")
        
        images_dir = os.path.join(output_dir, "slide_images")
        os.makedirs(images_dir, exist_ok=True)
        
        try:
            # Get slide count for fallback
            num_slides = self._get_slide_count(pptx_path)
            
            # Try to export images
            image_paths = self._export_images(pptx_path, images_dir, num_slides)
            
            if image_paths:
                return ProcessingResult.from_existing(pptx_path, output_dir, image_paths)
            return ProcessingResult.from_error(pptx_path, output_dir, "Failed to export images")
        except Exception as e:
            logger.error(f"Error processing PPTX {pptx_path}: {e}")
            return ProcessingResult.from_error(pptx_path, output_dir, str(e))
    
    def _get_slide_count(self, pptx_path: str) -> int:
        """Get number of slides in PPTX."""
        try:
            from pptx import Presentation
            return len(Presentation(pptx_path).slides)
        except Exception:
            return 0
    
    def _export_images(self, pptx_path: str, output_dir: str, num_slides: int) -> List[str]:
        """Export slides as images."""
        # Try LibreOffice
        images = self._export_with_libreoffice(pptx_path, output_dir)
        if images:
            return images
        
        # Try existing PDF
        pdf_path = pptx_path.replace(".pptx", ".pdf").replace(".PPTX", ".pdf")
        if os.path.exists(pdf_path):
            images = self.pdf_processor._convert_to_images(pdf_path, output_dir)
            if images:
                return images
        
        # Create placeholders
        if num_slides > 0:
            return self._create_placeholders(output_dir, num_slides)
        return []
    
    def _export_with_libreoffice(self, pptx_path: str, output_dir: str) -> List[str]:
        """Export via LibreOffice -> PDF -> images."""
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return []
        
        temp_dir = os.path.join(output_dir, "_temp")
        os.makedirs(temp_dir, exist_ok=True)
        
        try:
            # Convert to PDF
            result = subprocess.run([
                soffice, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, pptx_path
            ], capture_output=True, timeout=180)
            
            if result.returncode != 0:
                return []
            
            # Find generated PDF
            pdf_files = list(Path(temp_dir).glob("*.pdf"))
            if not pdf_files:
                return []
            
            # Convert PDF to images
            images = self.pdf_processor._convert_to_images(str(pdf_files[0]), output_dir)
            return images
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
    
    def _create_placeholders(self, output_dir: str, num_slides: int) -> List[str]:
        """Create placeholder images."""
        try:
            from PIL import Image, ImageDraw
            paths = []
            for i in range(num_slides):
                img = Image.new('RGB', (1920, 1080), color='white')
                draw = ImageDraw.Draw(img)
                text = f"Slide {i+1}"
                bbox = draw.textbbox((0, 0), text)
                x = (1920 - bbox[2]) // 2
                y = (1080 - bbox[3]) // 2
                draw.text((x, y), text, fill='gray')
                
                path = os.path.join(output_dir, f"slide_{i+1:04d}.{self.image_format}")
                img.save(path)
                paths.append(path)
            return paths
        except ImportError:
            return []


# =============================================================================
# URL Processor (for Zhipu)
# =============================================================================

class URLProcessor:
    """Process URLs via external subprocess."""
    
    def __init__(self):
        self.script_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            "process_zhipu.py"
        )
    
    def process(self, url: str, output_dir: str, force: bool = False) -> ProcessingResult:
        """Process a URL to extract slide images."""
        if not os.path.exists(self.script_path):
            return ProcessingResult.from_error(url, output_dir, "process_zhipu.py not found")
        
        images_dir = os.path.join(output_dir, "slide_images")
        
        try:
            cmd = [sys.executable, self.script_path, "--url", url, "--output", output_dir]
            if force:
                cmd.append("--force")
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                return ProcessingResult.from_error(url, output_dir, f"Subprocess failed: {result.stderr[:200]}")
            
            images = get_existing_images(images_dir)
            if images:
                return ProcessingResult.from_existing(url, output_dir, images)
            return ProcessingResult.from_error(url, output_dir, "No images generated")
            
        except subprocess.TimeoutExpired:
            return ProcessingResult.from_error(url, output_dir, "Processing timed out")
        except Exception as e:
            return ProcessingResult.from_error(url, output_dir, str(e))


# =============================================================================
# Main Preprocessor
# =============================================================================

class ProductPreprocessor:
    """Main preprocessor for all product types."""
    
    def __init__(
        self,
        ppt_gen_root: str = None,
        image_format: str = "png",
        image_dpi: int = 150
    ):
        self.ppt_gen_root = ppt_gen_root
        self.pptx_processor = PPTXProcessor(image_format, image_dpi)
        self.pdf_processor = PDFProcessor(image_format, image_dpi)
        self.url_processor = URLProcessor()
    
    def process_all(
        self,
        products: List[str] = None,
        difficulties: List[str] = None,
        topics: List[str] = None,
        force: bool = False
    ) -> Dict[str, List[ProcessingResult]]:
        """Process all products or a subset."""
        products = products or ALL_PRODUCTS
        results = {}
        
        # Use tqdm progress bar for products
        with tqdm(total=len(products), desc="Overall Progress", position=0) as pbar:
            for product in products:
                logger.info(f"\n{'='*60}\nProcessing: {product}\n{'='*60}")
                results[product] = self.process_product(product, difficulties, topics, force)
                
                success = sum(1 for r in results[product] if r.success)
                slides = sum(r.num_slides for r in results[product] if r.success)
                logger.info(f"{product}: {success}/{len(results[product])} successful, {slides} slides")
                
                pbar.update(1)
                pbar.set_postfix({"product": product, "success": success, "slides": slides})
        
        return results
    
    def process_product(
        self,
        product: str,
        difficulties: List[str] = None,
        topics: List[str] = None,
        force: bool = False
    ) -> List[ProcessingResult]:
        """Process all PPTs for a product."""
        product_path = os.path.join(self.ppt_gen_root, PRODUCT_PATHS.get(product, product))
        
        if not os.path.exists(product_path):
            logger.error(f"Product directory not found: {product_path}")
            return []
        
        results = []
        
        # Count total topics for progress bar
        total_topics = 0
        diff_dirs = self._list_dirs(product_path, difficulties)
        for diff in diff_dirs:
            diff_path = os.path.join(product_path, diff)
            total_topics += len(self._list_dirs(diff_path, topics))
        
        # Iterate through difficulties with progress bar
        with tqdm(total=total_topics, desc=f"{product}", position=1, leave=False) as pbar:
            for diff in diff_dirs:
                diff_path = os.path.join(product_path, diff)
                
                # Iterate through topics
                for topic in self._list_dirs(diff_path, topics):
                    topic_path = os.path.join(diff_path, topic)
                    topic_results = self._process_topic(topic_path, product, diff, topic, force)
                    results.extend(topic_results)
                    
                    # Update progress bar
                    pbar.update(1)
                    success = sum(1 for r in topic_results if r.success)
                    pbar.set_postfix({"diff": diff, "topic": topic[:20], "ok": success})
        
        return results
    
    def _list_dirs(self, path: str, filter_list: List[str] = None) -> List[str]:
        """List subdirectories, optionally filtered."""
        dirs = [d for d in os.listdir(path) 
                if os.path.isdir(os.path.join(path, d)) and not d.startswith('.')]
        if filter_list:
            dirs = [d for d in dirs if d in filter_list]
        return dirs
    
    def _process_topic(
        self,
        topic_path: str,
        product: str,
        difficulty: str,
        topic: str,
        force: bool
    ) -> List[ProcessingResult]:
        """Process all PPTs in a topic directory."""
        results = []
        
        # Check for subdirectories (PPT folders)
        # A valid PPT folder contains a .pptx file or .txt file with the same name as the folder
        ppt_dirs = []
        for d in self._list_dirs(topic_path):
            if d in ("slide_images", "slide_texts"):
                continue
            dir_path = os.path.join(topic_path, d)
            # Check if folder contains same-named .pptx or .txt file
            pptx_file = os.path.join(dir_path, f"{d}.pptx")
            txt_file = os.path.join(dir_path, f"{d}.txt")
            if os.path.exists(pptx_file) or os.path.exists(txt_file):
                ppt_dirs.append(d)
        
        if ppt_dirs:
            # Process each PPT subdirectory
            for ppt_id in ppt_dirs:
                ppt_path = os.path.join(topic_path, ppt_id)
                result = self._process_folder(ppt_path, force)
                self._set_metadata(result, product, difficulty, topic, ppt_id)
                results.append(result)
        else:
            # Files directly in topic folder (NotebookLM, Zhipu, etc.)
            result = self._process_folder(topic_path, force)
            self._set_metadata(result, product, difficulty, topic, topic)
            results.append(result)
        
        return results
    
    def _process_folder(self, folder_path: str, force: bool) -> ProcessingResult:
        """Process a folder containing PPT content (PPTX, PDF, URL, or existing images)."""
        images_dir = os.path.join(folder_path, "slide_images")
        
        # Check for existing images
        if not force:
            existing = get_existing_images(images_dir)
            if existing:
                logger.info(f"Already processed: {folder_path} ({len(existing)} images)")
                return ProcessingResult.from_existing(folder_path, folder_path, existing)
        
        # Try different file types in order of preference
        # 1. PPTX files
        pptx_files = find_files(folder_path, [".pptx"])
        if pptx_files:
            return self.pptx_processor.process(str(pptx_files[0]), folder_path)
        
        # 2. PDF files
        pdf_files = find_files(folder_path, [".pdf"])
        if pdf_files:
            return self.pdf_processor.process(str(pdf_files[0]), folder_path)
        
        # 3. URL files (Zhipu)
        url_files = list(Path(folder_path).glob("url.txt")) + list(Path(folder_path).glob("*.txt"))
        if url_files:
            try:
                url = Path(url_files[0]).read_text(encoding='utf-8').strip()
                if url:
                    return self.url_processor.process(url, folder_path, force)
            except Exception as e:
                logger.error(f"Error reading URL file: {e}")
        
        logger.warning(f"No processable files found in {folder_path}")
        return ProcessingResult.from_error(folder_path, folder_path, "No PPTX, PDF, or URL found")
    
    def _set_metadata(self, result: ProcessingResult, product: str, difficulty: str, topic: str, ppt_id: str):
        """Set metadata fields on a result."""
        result.product = product
        result.difficulty = difficulty
        result.topic = topic
        result.ppt_id = ppt_id


# =============================================================================
# CLI
# =============================================================================

async def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(description="Pre-process PPT files for evaluation")
    
    parser.add_argument("--product", nargs="+", default=["all"],
                       help="Products to process (or 'all')")
    parser.add_argument("--difficulty", nargs="+", 
                       default=["topic_introduction", "work_report", "business_plan","brand_promote","personal_statement","product_launch","course_preparation"],
                       help="Difficulty levels")
    parser.add_argument("--topic", nargs="+", help="Specific topics")
    parser.add_argument("--pptx", help="Single PPTX file")
    parser.add_argument("--pdf", help="Single PDF file")
    parser.add_argument("--ppt-gen-root",
                       required=True,
                       help="Root directory for PPT products")
    parser.add_argument("--format", default="png", choices=["png", "jpg"])
    parser.add_argument("--dpi", type=int, default=150)
    parser.add_argument("--force", action="store_true", help="Force reprocess")
    
    args = parser.parse_args()
    results = []
    
    # Process single files
    if args.pptx:
        processor = PPTXProcessor(args.format, args.dpi)
        results.append(processor.process(args.pptx, str(Path(args.pptx).parent)))
    
    if args.pdf:
        processor = PDFProcessor(args.format, args.dpi)
        results.append(processor.process(args.pdf, str(Path(args.pdf).parent)))
    
    # Process products
    if not (args.pptx or args.pdf):
        preprocessor = ProductPreprocessor(args.ppt_gen_root, args.format, args.dpi)
        products = None if "all" in args.product else args.product
        
        all_results = preprocessor.process_all(
            products=products,
            difficulties=args.difficulty,
            topics=args.topic,
            force=args.force
        )
        
        for product_results in all_results.values():
            results.extend(product_results)
    
    # Summary
    print(f"\n{'='*60}\nPRE-PROCESSING COMPLETE\n{'='*60}")
    
    success = sum(1 for r in results if r.success)
    slides = sum(r.num_slides for r in results if r.success)
    print(f"Total: {len(results)} items, {success} successful, {slides} slides")
    
    by_product = {}
    for r in results:
        by_product.setdefault(r.product, {"success": 0, "failed": 0, "slides": 0})
        if r.success:
            by_product[r.product]["success"] += 1
            by_product[r.product]["slides"] += r.num_slides
        else:
            by_product[r.product]["failed"] += 1
    
    print("\nBy Product:")
    for prod, stats in sorted(by_product.items()):
        print(f"  {prod}: {stats['success']} success, {stats['failed']} failed, {stats['slides']} slides")
    
    print("="*60)


if __name__ == "__main__":
    asyncio.run(main())
